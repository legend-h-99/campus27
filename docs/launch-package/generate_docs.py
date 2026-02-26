#!/usr/bin/env python3
"""Generate Word, PowerPoint, and PDF documents for Campus27 launch package."""

import os
from pathlib import Path

# ─── Output directory ───
OUT = Path(__file__).parent
os.makedirs(OUT, exist_ok=True)

# ═══════════════════════════════════════════════════════════════
# 1. WORD DOCUMENTS (DOCX)
# ═══════════════════════════════════════════════════════════════
from docx import Document
from docx.shared import Inches, Pt, Cm, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.enum.table import WD_TABLE_ALIGNMENT

PRIMARY = RGBColor(0x2D, 0x7A, 0x8F)
SECONDARY = RGBColor(0x1E, 0x5A, 0x6B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
DARK = RGBColor(0x33, 0x33, 0x33)


def style_doc(doc, title_text, subtitle_text=""):
    """Apply consistent styling to a Word document."""
    # Default font
    style = doc.styles["Normal"]
    font = style.font
    font.name = "Calibri"
    font.size = Pt(11)
    font.color.rgb = DARK

    # Title
    title = doc.add_heading(level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = title.add_run(title_text)
    run.font.size = Pt(28)
    run.font.color.rgb = PRIMARY
    run.font.bold = True

    if subtitle_text:
        sub = doc.add_paragraph()
        sub.alignment = WD_ALIGN_PARAGRAPH.CENTER
        run = sub.add_run(subtitle_text)
        run.font.size = Pt(14)
        run.font.color.rgb = SECONDARY
        run.font.italic = True

    doc.add_paragraph()  # spacer
    return doc


def add_heading(doc, text, level=1):
    h = doc.add_heading(text, level=level)
    for run in h.runs:
        run.font.color.rgb = PRIMARY if level == 1 else SECONDARY
    return h


def add_table(doc, headers, rows):
    table = doc.add_table(rows=1 + len(rows), cols=len(headers))
    table.style = "Light Grid Accent 1"
    table.alignment = WD_TABLE_ALIGNMENT.CENTER

    # Header row
    for i, h in enumerate(headers):
        cell = table.rows[0].cells[i]
        cell.text = h
        for p in cell.paragraphs:
            p.alignment = WD_ALIGN_PARAGRAPH.CENTER
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(10)

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, val in enumerate(row):
            cell = table.rows[r_idx + 1].cells[c_idx]
            cell.text = str(val)
            for p in cell.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(10)

    doc.add_paragraph()  # spacer
    return table


# ─── BRD Document ───
def generate_brd():
    doc = Document()
    style_doc(doc, "Campus27", "وثيقة متطلبات الأعمال — Business Requirements Document")

    doc.add_paragraph("التاريخ: فبراير 2026  |  الإصدار: 1.0  |  الحالة: مسودة للمراجعة").alignment = WD_ALIGN_PARAGRAPH.CENTER

    # Executive Summary
    add_heading(doc, "1. الملخص التنفيذي")
    doc.add_paragraph(
        "منصة Campus27 هي نظام متكامل لإدارة الكليات التقنية، مصمم خصيصاً "
        "للمؤسسة العامة للتدريب التقني والمهني (TVTC). تجمع المنصة بين إدارة الشؤون "
        "الأكاديمية والمالية والموارد البشرية وضمان الجودة في واجهة موحدة ثنائية اللغة "
        "(عربي/إنجليزي) مع دعم كامل للاتجاه من اليمين لليسار."
    )
    doc.add_paragraph(
        "الهدف الرئيسي: رقمنة وأتمتة العمليات الإدارية والأكاديمية في الكلية التقنية، "
        "مع تحقيق متطلبات الاعتماد الأكاديمي (NCAAA) وتوفير رؤى ذكية مدعومة بالذكاء الاصطناعي."
    )

    # Scope
    add_heading(doc, "2. نطاق المنصة")
    add_heading(doc, "2.1 الوحدات الرئيسية (11 وحدة)", level=2)
    add_table(doc,
        ["#", "الوحدة", "الوصف"],
        [
            ["1", "لوحة القيادة", "لوحات ذكية حسب الدور (عميد/مدرب/متدرب)"],
            ["2", "شؤون المتدربين", "تسجيل، متابعة، حالات، مستويات"],
            ["3", "شؤون المدربين", "ملفات، تخصصات، أقسام، تقييمات"],
            ["4", "الأقسام والمقررات", "إدارة الأقسام، المقررات، الجداول، الدرجات"],
            ["5", "الحضور والانصراف", "تسجيل يومي، تقارير، إنذارات"],
            ["6", "الشؤون المالية", "ميزانيات، معاملات، مشتريات، مستودعات"],
            ["7", "ضمان الجودة", "معايير NCAAA، مؤشرات أداء، تدقيق، اعتماد"],
            ["8", "الموارد البشرية", "إجازات، تقييمات أداء، رواتب"],
            ["9", "التعلم الإلكتروني", "محتوى رقمي، اختبارات إلكترونية"],
            ["10", "إدارة المهام والمشاريع", "تتبع مهام، مشاريع، ميزانيات"],
            ["11", "مركز الذكاء الاصطناعي", "محادثة ذكية، تنبؤات، إنذار مبكر، تقارير آلية"],
        ]
    )

    # Roles
    add_heading(doc, "2.2 الأدوار والصلاحيات (13 دور)", level=2)
    add_table(doc,
        ["الدور", "النطاق"],
        [
            ["مدير النظام", "صلاحيات كاملة"],
            ["عميد الكلية", "صلاحيات كاملة + لوحة قيادة تنفيذية"],
            ["وكيل شؤون المدربين", "المدربون، الجداول، الدرجات، المهام"],
            ["وكيل شؤون المتدربين", "المتدربون، الدرجات، الحضور، التعلم الإلكتروني"],
            ["وكيل الجودة", "نظام الجودة الكامل، الاعتماد"],
            ["رئيس القسم", "قسمه، مدربوه، مقرراته"],
            ["المدرب", "مقرراته، الدرجات، الحضور"],
            ["المتدرب", "درجاته، جدوله، التعلم الإلكتروني"],
            ["المحاسب", "الشؤون المالية، التقارير"],
            ["مسؤول الجودة", "مؤشرات الأداء، القياس"],
            ["مدير الموارد البشرية", "الإجازات، التقييمات"],
            ["مسؤول تقنية المعلومات", "إدارة المستخدمين، الإعدادات"],
            ["منسق الوحدة", "المهام، التقارير"],
        ]
    )

    # Technical capabilities
    add_heading(doc, "2.3 القدرات التقنية", level=2)
    add_table(doc,
        ["الميزة", "التفاصيل"],
        [
            ["ثنائي اللغة", "عربي (RTL) + إنجليزي (LTR)"],
            ["الذكاء الاصطناعي", "5 خدمات: محادثة، رؤى، إنذار مبكر، توصيات، تقارير آلية"],
            ["قاعدة البيانات", "44 جدول يغطي جميع العمليات"],
            ["الأمان", "مصادقة JWT، صلاحيات دقيقة (59 صلاحية)"],
            ["الاستضافة", "سحابية (Vercel + Neon PostgreSQL)"],
            ["التقارير", "PDF + Excel لجميع الوحدات"],
        ]
    )

    # Functional Requirements
    add_heading(doc, "3. المتطلبات الوظيفية الرئيسية")

    add_heading(doc, "3.1 لوحة القيادة الذكية", level=2)
    for req in [
        "FR-01: عرض إحصائيات فورية (متدربين، مدربين، أقسام، مقررات)",
        "FR-02: رسوم بيانية تفاعلية (توزيع المتدربين، معدلات الحضور)",
        "FR-03: لوحة مخصصة لكل دور",
        "FR-04: رؤى الذكاء الاصطناعي (تنبؤات، توصيات، تنبيهات)",
    ]:
        doc.add_paragraph(req, style="List Bullet")

    add_heading(doc, "3.2 إدارة الشؤون الأكاديمية", level=2)
    for req in [
        "FR-05: تسجيل ومتابعة 300+ متدرب عبر 4 مستويات",
        "FR-06: إدارة 120 مقرر عبر 10 أقسام",
        "FR-07: جدولة أكاديمية (4 فترات × 5 أيام)",
        "FR-08: نظام درجات متكامل مع سير عمل اعتماد",
        "FR-09: تسجيل حضور يومي مع تقارير وإنذارات",
    ]:
        doc.add_paragraph(req, style="List Bullet")

    add_heading(doc, "3.3 ضمان الجودة والاعتماد (NCAAA)", level=2)
    for req in [
        "FR-10: 10 معايير جودة متوافقة مع NCAAA v3.0",
        "FR-11: 15+ مؤشر أداء رئيسي مع أهداف وقياسات دورية",
        "FR-12: نظام تدقيق (داخلي، خارجي، مراجعة برامج، تقييم ذاتي)",
        "FR-13: خطط تحسين مع إجراءات تصحيحية/وقائية",
        "FR-14: استبيانات (7 أنواع)",
        "FR-15: تحليل جاهزية الاعتماد بالذكاء الاصطناعي",
    ]:
        doc.add_paragraph(req, style="List Bullet")

    add_heading(doc, "3.4 الشؤون المالية", level=2)
    for req in [
        "FR-16: إدارة ميزانية سنوية (32.5 مليون ريال عبر 12 بند)",
        "FR-17: معاملات مالية مع سير عمل اعتماد",
        "FR-18: طلبات مشتريات ومستودعات",
    ]:
        doc.add_paragraph(req, style="List Bullet")

    add_heading(doc, "3.5 الذكاء الاصطناعي", level=2)
    for req in [
        "FR-19: محادثة ذكية للاستفسارات الإدارية",
        "FR-20: نظام إنذار مبكر للمتدربين المعرضين للخطر",
        "FR-21: تنبؤات أداء بنسبة دقة 86%",
        "FR-22: توصيات ذكية قابلة للتنفيذ",
        "FR-23: تقارير آلية مولدة بالذكاء الاصطناعي",
    ]:
        doc.add_paragraph(req, style="List Bullet")

    # Non-functional
    add_heading(doc, "4. المتطلبات غير الوظيفية")
    add_table(doc,
        ["المتطلب", "المعيار"],
        [
            ["الأداء", "تحميل الصفحات < 3 ثوان"],
            ["التوفر", "99.5% (استضافة سحابية)"],
            ["الأمان", "تشفير JWT، HTTPS، صلاحيات على مستوى العمليات"],
            ["التوسع", "بنية SaaS قابلة لخدمة كليات متعددة"],
            ["التوافق", "متصفحات حديثة + أجهزة محمولة (متجاوب)"],
            ["اللغات", "عربي + إنجليزي مع تبديل فوري"],
        ]
    )

    # Current status
    add_heading(doc, "5. الوضع الحالي")
    add_table(doc,
        ["البند", "الحالة"],
        [
            ["التطوير", "مكتمل (النسخة 1.0)"],
            ["الاستضافة", "مفعلة على Vercel"],
            ["قاعدة البيانات", "مفعلة على Neon PostgreSQL"],
            ["البيانات التجريبية", "محملة (377 مستخدم، 120 مقرر، 10 أقسام)"],
            ["الذكاء الاصطناعي", "مفعل ويعمل"],
            ["الرابط", "https://campus27.vercel.app"],
            ["بيانات الدخول التجريبية", "admin@campus27.sa / 123456"],
        ]
    )

    # Success criteria
    add_heading(doc, "6. معايير النجاح")
    for c in [
        "اعتماد 80% من المستخدمين المستهدفين خلال 30 يوم من الإطلاق التجريبي",
        "تحسين كفاءة العمليات الإدارية بنسبة 40%",
        "تجهيز 100% من وثائق الجودة المطلوبة للاعتماد",
        "رضا المستخدمين ≥ 4/5",
    ]:
        doc.add_paragraph(c, style="List Number")

    # Footer
    doc.add_paragraph()
    p = doc.add_paragraph("إعداد: فريق التطوير التقني | Campus27 v1.0 | فبراير 2026")
    p.alignment = WD_ALIGN_PARAGRAPH.CENTER

    doc.save(OUT / "Campus27-BRD.docx")
    print("✓ Campus27-BRD.docx")


# ─── Pilot Launch Plan Document ───
def generate_pilot_plan():
    doc = Document()
    style_doc(doc, "Campus27", "خطة إطلاق التجربة المحدودة — Pilot Launch Plan")

    doc.add_paragraph("المدة: 8 أسابيع  |  المشاركون: ~30 شخص").alignment = WD_ALIGN_PARAGRAPH.CENTER

    add_heading(doc, "1. أهداف التجربة")
    add_table(doc,
        ["#", "الهدف", "مؤشر القياس", "المستهدف"],
        [
            ["1", "التحقق من ملاءمة المنصة", "نسبة إكمال المهام", "≥ 80%"],
            ["2", "قياس سهولة الاستخدام", "تقييم رضا المستخدمين", "≥ 4/5"],
            ["3", "اختبار نظام الجودة", "مؤشرات أداء مسجلة", "≥ 10 KPIs"],
            ["4", "تقييم فائدة الذكاء الاصطناعي", "دقة التنبؤات", "≥ 80%"],
            ["5", "جمع الملاحظات للتحسين", "ملاحظات معالجة", "≥ 90%"],
        ]
    )

    add_heading(doc, "2. المشاركون")
    add_table(doc,
        ["الفئة", "العدد", "معيار الاختيار"],
        [
            ["وكلاء الكلية", "3", "وكيل المدربين + المتدربين + الجودة"],
            ["رؤساء الأقسام", "3", "أقسام متنوعة"],
            ["مدربون", "20", "6-7 من كل قسم مشارك"],
            ["مسؤول الجودة", "1", "اختبار وحدة الجودة"],
            ["المحاسب", "1", "اختبار الوحدة المالية"],
            ["مسؤول IT + منسق", "2", "الدعم والتنسيق"],
        ]
    )

    add_heading(doc, "3. الجدول الزمني")

    add_heading(doc, "المرحلة 1: التهيئة (الأسبوع 1)", level=2)
    for task in [
        "إنشاء حسابات المستخدمين الحقيقية",
        "إدخال بيانات الأقسام والمقررات الفعلية",
        "إدخال بيانات المدربين المشاركين",
        "اختبار تقني شامل",
    ]:
        doc.add_paragraph(task, style="List Bullet")

    add_heading(doc, "المرحلة 2: التدريب (الأسبوع 2)", level=2)
    add_table(doc,
        ["الجلسة", "الفئة", "المدة", "المحتوى"],
        [
            ["1", "وكلاء الكلية", "2 ساعة", "لوحة القيادة، التقارير، AI"],
            ["2", "رؤساء الأقسام", "3 ساعات", "إدارة القسم، المدربون، المقررات"],
            ["3-4", "المدربون", "3 ساعات × 2", "الدرجات، الحضور، المهام"],
            ["5", "مسؤول الجودة", "2 ساعة", "نظام الجودة الكامل"],
        ]
    )

    add_heading(doc, "المرحلة 3: التجربة الفعلية (الأسابيع 3-6)", level=2)
    add_table(doc,
        ["الأسبوع", "التركيز", "المهام"],
        [
            ["3", "البدء الناعم", "حضور يومي + درجات اختبار واحد"],
            ["4", "التوسع", "المهام + المالية + مؤشرات الجودة"],
            ["5", "الاستقرار", "جميع العمليات اليومية عبر المنصة"],
            ["6", "التقييم الأولي", "مراجعة التقارير + تقييم AI"],
        ]
    )

    add_heading(doc, "المرحلة 4: التقييم والقرار (الأسابيع 7-8)", level=2)
    for task in [
        "استبيان رضا المستخدمين",
        "جمع وتحليل الملاحظات",
        "تقرير نتائج التجربة",
        "عرض النتائج على العميد",
        "قرار التوسع أو التعديل",
    ]:
        doc.add_paragraph(task, style="List Bullet")

    add_heading(doc, "4. معايير نجاح التجربة")
    add_table(doc,
        ["المعيار", "الحد الأدنى", "المثالي"],
        [
            ["نسبة الاستخدام اليومي", "60%", "80%+"],
            ["رضا المستخدمين", "3.5/5", "4+/5"],
            ["المشاكل التقنية الحرجة", "≤ 3", "0"],
            ["وقت إنجاز المهام", "أسرع من الحالي", "أسرع بـ 50%+"],
        ]
    )

    add_heading(doc, "5. المخاطر والتخفيف")
    add_table(doc,
        ["المخاطر", "الاحتمال", "التخفيف"],
        [
            ["مقاومة التغيير", "متوسط", "تدريب عملي + دعم مستمر"],
            ["مشاكل تقنية", "منخفض", "دعم فني فوري + خطة طوارئ"],
            ["ضعف إدخال البيانات", "متوسط", "متابعة يومية + تبسيط واجهات"],
            ["انشغال المشاركين", "عالي", "تقليل المهام + مرونة في الجدول"],
        ]
    )

    add_heading(doc, "6. المطلوب للبدء")
    doc.add_paragraph("من العميد:", style="List Bullet")
    for item in ["الموافقة على خطة التجربة", "تحديد 3 أقسام مشاركة", "تعيين منسق مشروع"]:
        p = doc.add_paragraph(item)
        p.style = "List Bullet 2" if "List Bullet 2" in [s.name for s in doc.styles] else "List Bullet"

    doc.add_paragraph()
    p = doc.add_paragraph("تاريخ البدء المقترح: خلال أسبوعين من الموافقة")
    p.runs[0].bold = True

    doc.save(OUT / "Campus27-Pilot-Plan.docx")
    print("✓ Campus27-Pilot-Plan.docx")


# ═══════════════════════════════════════════════════════════════
# 2. POWERPOINT PRESENTATION (PPTX)
# ═══════════════════════════════════════════════════════════════
from pptx import Presentation
from pptx.util import Inches, Pt as PptPt, Emu
from pptx.dml.color import RGBColor as PptRGB
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


TEAL = PptRGB(0x2D, 0x7A, 0x8F)
DARK_TEAL = PptRGB(0x1E, 0x5A, 0x6B)
ACCENT = PptRGB(0x32, 0xB8, 0xC6)
PPT_WHITE = PptRGB(0xFF, 0xFF, 0xFF)
PPT_DARK = PptRGB(0x33, 0x33, 0x33)
LIGHT_BG = PptRGB(0xF0, 0xF8, 0xFA)


def add_slide(prs, layout_idx=1):
    """Add a slide using the given layout index."""
    layout = prs.slide_layouts[min(layout_idx, len(prs.slide_layouts) - 1)]
    return prs.slides.add_slide(layout)


def set_bg(slide, color):
    """Set slide background color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18, bold=False, color=PPT_DARK, alignment=PP_ALIGN.RIGHT):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = PptPt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = alignment
    return txBox


def add_bullet_slide(prs, title_text, bullets, bg_color=PPT_WHITE):
    slide = add_slide(prs, 1)
    set_bg(slide, bg_color)

    # Title
    add_textbox(slide, 0.5, 0.3, 9, 0.8, title_text, font_size=28, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)

    # Bullets
    txBox = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.5))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, bullet in enumerate(bullets):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = f"  {bullet}"
        p.font.size = PptPt(16)
        p.font.color.rgb = PPT_DARK
        p.space_after = PptPt(12)
        p.alignment = PP_ALIGN.RIGHT

    return slide


def add_table_slide(prs, title_text, headers, rows, bg_color=PPT_WHITE):
    slide = add_slide(prs, 1)
    set_bg(slide, bg_color)

    add_textbox(slide, 0.5, 0.2, 9, 0.7, title_text, font_size=26, bold=True, color=TEAL, alignment=PP_ALIGN.CENTER)

    num_rows = len(rows) + 1
    num_cols = len(headers)
    tbl_width = Inches(9)
    tbl_height = Inches(min(5.5, 0.5 * num_rows))

    shape = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.2), tbl_width, tbl_height)
    table = shape.table

    # Header
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = PptPt(12)
            p.font.bold = True
            p.font.color.rgb = PPT_WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = TEAL

    # Rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = PptPt(11)
                p.alignment = PP_ALIGN.CENTER
            if r % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_BG

    return slide


def generate_pptx():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # ─── Slide 1: Cover ───
    slide = add_slide(prs, 0)
    set_bg(slide, TEAL)
    add_textbox(slide, 1, 1.5, 8, 1.5, "Campus27", font_size=48, bold=True, color=PPT_WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 3, 8, 1, "نظام إدارة الكلية التقنية المتكامل", font_size=24, color=PPT_WHITE, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 4.5, 8, 0.6, "عرض تقديمي — فبراير 2026", font_size=16, color=ACCENT, alignment=PP_ALIGN.CENTER)

    # ─── Slide 2: Challenge ───
    add_bullet_slide(prs, "التحدي — الوضع الحالي", [
        "أنظمة متعددة ومنفصلة لإدارة العمليات",
        "صعوبة الحصول على صورة شاملة عن أداء الكلية",
        "العمليات الورقية تستهلك وقت الكادر الإداري",
        "متطلبات الاعتماد (NCAAA) تحتاج توثيق مستمر",
        "غياب نظام إنذار مبكر للمتدربين المعرضين للتعثر",
    ], bg_color=LIGHT_BG)

    # ─── Slide 3: Solution ───
    add_bullet_slide(prs, "الحل — Campus27 كل شيء في مكان واحد", [
        "11 وحدة إدارية متكاملة في منصة واحدة",
        "لوحة قيادة ذكية مخصصة لكل دور",
        "نظام جودة كامل متوافق مع NCAAA v3.0",
        "ذكاء اصطناعي مدمج (5 خدمات)",
        "ثنائي اللغة: عربي + إنجليزي مع تبديل فوري",
        "استضافة سحابية — أي جهاز، أي مكان",
    ])

    # ─── Slide 4: Numbers ───
    add_table_slide(prs, "المنصة بالأرقام",
        ["المؤشر", "القيمة"],
        [
            ["وحدات إدارية", "11 وحدة"],
            ["صفحات المنصة", "22 صفحة"],
            ["أدوار المستخدمين", "13 دور"],
            ["الصلاحيات", "59 صلاحية دقيقة"],
            ["جداول البيانات", "44 جدول"],
            ["نقاط API", "39+ نقطة وصول"],
            ["خدمات AI", "5 خدمات"],
            ["اللغات", "عربي + إنجليزي"],
        ],
        bg_color=LIGHT_BG
    )

    # ─── Slide 5: Dashboard ───
    add_bullet_slide(prs, "لوحة القيادة الذكية", [
        "لوحة العميد/المدير — رؤية شاملة لجميع مؤشرات الكلية",
        "لوحة المدرب — مقرراته، درجاته، حضوره",
        "لوحة المتدرب — درجاته، جدوله، إشعاراته",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "إحصائيات فورية + رسوم بيانية تفاعلية",
        "رؤى الذكاء الاصطناعي + الموافقات المعلقة",
    ])

    # ─── Slide 6: Quality ───
    add_table_slide(prs, "نظام الجودة والاعتماد — NCAAA v3.0",
        ["القدرة", "التفاصيل"],
        [
            ["المعايير", "10 معايير NCAAA"],
            ["مؤشرات الأداء", "15+ KPI مع أهداف وقياسات"],
            ["التدقيق", "داخلي، خارجي، مراجعة برامج"],
            ["خطط التحسين", "إجراءات تصحيحية ووقائية"],
            ["الاستبيانات", "7 أنواع"],
            ["الوثائق", "سياسات، إجراءات، أدلة، نماذج"],
            ["AI", "تحليل جاهزية الاعتماد تلقائياً"],
        ]
    )

    # ─── Slide 7: AI ───
    add_bullet_slide(prs, "مركز الذكاء الاصطناعي — 5 خدمات ذكية", [
        "1. المحادثة الذكية — اسأل أي سؤال عن بيانات الكلية",
        "2. الرؤى والتحليلات — اكتشاف أنماط وفرص تحسين",
        "3. الإنذار المبكر — المتدربين المعرضين للتعثر (دقة 86%)",
        "4. التوصيات الذكية — اقتراحات قابلة للتنفيذ",
        "5. التقارير الآلية — توليد تقارير شاملة تلقائياً",
    ], bg_color=LIGHT_BG)

    # ─── Slide 8: Security ───
    add_bullet_slide(prs, "الأمان والصلاحيات", [
        "13 دور مستخدم بـ 59 صلاحية دقيقة",
        "كل مستخدم يرى فقط ما يخصه — لا أكثر",
        "تشفير JWT + اتصالات HTTPS مشفرة",
        "سجل مراجعة كامل لجميع العمليات (Audit Log)",
        "━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━",
        "عميد ← يرى الكل | وكيل ← نطاقه | مدرب ← مقرراته | متدرب ← بياناته",
    ])

    # ─── Slide 9: Pilot Plan ───
    add_table_slide(prs, "خطة الإطلاق التجريبي — 8 أسابيع",
        ["المرحلة", "المدة", "المشاركون", "الهدف"],
        [
            ["التهيئة", "أسبوع 1", "فريق تقني", "إعداد البيئة والبيانات"],
            ["التدريب", "أسبوع 2", "20 مدرب + 3 رؤساء + وكلاء", "ورش تدريبية"],
            ["التجربة", "أسابيع 3-6", "جميع المشاركين", "استخدام فعلي"],
            ["التقييم", "أسابيع 7-8", "فريق المشروع", "تقييم واتخاذ قرار"],
        ],
        bg_color=LIGHT_BG
    )

    # ─── Slide 10: Next Steps ───
    slide = add_slide(prs, 0)
    set_bg(slide, TEAL)
    add_textbox(slide, 1, 0.8, 8, 1, "الخطوة التالية", font_size=36, bold=True, color=PPT_WHITE, alignment=PP_ALIGN.CENTER)

    steps = [
        "1. الموافقة على إطلاق التجربة المحدودة",
        "2. تحديد 3 أقسام للمشاركة في التجربة",
        "3. تعيين منسق مشروع من جانب الكلية",
        "4. تخصيص فترة تدريبية (يومين) للمشاركين",
    ]
    txBox = slide.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(7), Inches(3))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, step in enumerate(steps):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.text = step
        p.font.size = PptPt(20)
        p.font.color.rgb = PPT_WHITE
        p.space_after = PptPt(16)
        p.alignment = PP_ALIGN.RIGHT

    add_textbox(slide, 1, 5.5, 8, 0.5, "بداية التجربة: خلال أسبوعين من الموافقة", font_size=18, bold=True, color=ACCENT, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1, 6.2, 8, 0.8, "campus27.vercel.app\nadmin@campus27.sa / 123456", font_size=14, color=PPT_WHITE, alignment=PP_ALIGN.CENTER)

    prs.save(OUT / "Campus27-Presentation.pptx")
    print("✓ Campus27-Presentation.pptx")


# ═══════════════════════════════════════════════════════════════
# 3. PDF DOCUMENTS
# ═══════════════════════════════════════════════════════════════
from fpdf import FPDF
import os

# Find an Arabic-supporting font
FONT_PATH = None
FONT_CANDIDATES = [
    "/System/Library/Fonts/Supplemental/Arial Unicode.ttf",
    "/System/Library/Fonts/Supplemental/Tahoma.ttf",
    "/Library/Fonts/Arial Unicode.ttf",
    "/System/Library/Fonts/GeezaPro.ttc",
    "/System/Library/Fonts/Supplemental/Al Nile.ttc",
]

for fp in FONT_CANDIDATES:
    if os.path.exists(fp):
        FONT_PATH = fp
        break


class ArabicPDF(FPDF):
    """PDF with Unicode font support."""

    def __init__(self):
        super().__init__()
        self.set_auto_page_break(auto=True, margin=20)
        # Use Arial Unicode MS which supports full Unicode
        font_path = "/System/Library/Fonts/Supplemental/Arial Unicode.ttf"
        self.add_font("ArialUni", "", font_path)
        self.add_font("ArialUni", "B", font_path)
        self.add_font("ArialUni", "I", font_path)
        self.default_font = "ArialUni"

    def header_block(self, title, subtitle=""):
        self.set_font(self.default_font, "B", 24)
        self.set_text_color(45, 122, 143)  # PRIMARY
        self.cell(0, 15, title, new_x="LMARGIN", new_y="NEXT", align="C")
        if subtitle:
            self.set_font(self.default_font, "", 12)
            self.set_text_color(30, 90, 107)  # SECONDARY
            self.cell(0, 8, subtitle, new_x="LMARGIN", new_y="NEXT", align="C")
        self.ln(10)

    def section_heading(self, text):
        self.ln(5)
        self.set_font(self.default_font, "B", 16)
        self.set_text_color(45, 122, 143)
        self.cell(0, 10, text, new_x="LMARGIN", new_y="NEXT")
        self.set_draw_color(45, 122, 143)
        self.line(10, self.get_y(), 200, self.get_y())
        self.ln(3)

    def body_text(self, text):
        self.set_font(self.default_font, "", 11)
        self.set_text_color(51, 51, 51)
        w = self.w - self.l_margin - self.r_margin
        self.multi_cell(w, 6, text)
        self.ln(3)

    def bullet(self, text):
        self.set_font(self.default_font, "", 11)
        self.set_text_color(51, 51, 51)
        w = self.w - self.l_margin - self.r_margin
        self.multi_cell(w, 6, "  - " + text)

    def simple_table(self, headers, rows):
        col_width = (self.w - 20) / len(headers)

        # Header
        self.set_font(self.default_font, "B", 10)
        self.set_fill_color(45, 122, 143)
        self.set_text_color(255, 255, 255)
        for h in headers:
            self.cell(col_width, 8, h, border=1, fill=True, align="C")
        self.ln()

        # Rows
        self.set_font(self.default_font, "", 9)
        self.set_text_color(51, 51, 51)
        for i, row in enumerate(rows):
            if i % 2 == 0:
                self.set_fill_color(240, 248, 250)
            else:
                self.set_fill_color(255, 255, 255)
            for val in row:
                self.cell(col_width, 7, str(val), border=1, fill=True, align="C")
            self.ln()
        self.ln(5)


def generate_brd_pdf():
    pdf = ArabicPDF()
    pdf.add_page()

    pdf.header_block("Campus27", "Business Requirements Document (BRD)")
    pdf.body_text("Version 1.0 | February 2026 | Draft for Review")

    pdf.section_heading("1. Executive Summary")
    pdf.body_text(
        "Campus27 is an integrated management platform for technical colleges, designed for TVTC. "
        "It combines academic, financial, HR, and quality management in a unified bilingual interface "
        "(Arabic/English) with full RTL support and AI-powered insights."
    )

    pdf.section_heading("2. Platform Scope")
    pdf.body_text("11 Management Modules:")
    pdf.simple_table(
        ["#", "Module", "Description"],
        [
            ["1", "Dashboard", "Role-based smart dashboards (Dean/Trainer/Trainee)"],
            ["2", "Trainees", "Registration, tracking, status, levels"],
            ["3", "Trainers", "Profiles, specializations, evaluations"],
            ["4", "Academics", "Departments, courses, schedules, grades"],
            ["5", "Attendance", "Daily recording, reports, alerts"],
            ["6", "Finance", "Budgets, transactions, procurement, warehouse"],
            ["7", "Quality (NCAAA)", "Standards, KPIs, audits, accreditation"],
            ["8", "HR", "Leaves, evaluations, payroll"],
            ["9", "E-Learning", "Digital content, online exams"],
            ["10", "Tasks & Projects", "Task tracking, project management"],
            ["11", "AI Center", "Chat, insights, early warning, auto-reports"],
        ]
    )

    pdf.section_heading("3. Roles & Permissions")
    pdf.body_text("13 user roles with 59 granular permissions ensuring each user sees only what they need.")
    pdf.simple_table(
        ["Role", "Scope"],
        [
            ["Super Admin / Dean", "Full access + executive dashboard"],
            ["VP Trainers", "Trainers, schedules, grades, tasks"],
            ["VP Trainees", "Trainees, grades, attendance, e-learning"],
            ["VP Quality", "Full quality system, accreditation"],
            ["Dept Head", "Own department, trainers, courses"],
            ["Trainer", "Own courses, grades, attendance"],
            ["Trainee", "Own grades, schedule, e-learning"],
            ["Accountant", "Finance module, reports"],
            ["Quality Officer", "KPIs, measurements"],
            ["HR Manager", "Leaves, evaluations"],
            ["IT Admin", "User management, settings"],
            ["Unit Coordinator", "Tasks, reports"],
        ]
    )

    pdf.section_heading("4. AI Capabilities")
    pdf.bullet("AI Chat — Ask any question about college data")
    pdf.bullet("Insights & Analytics — Pattern discovery and improvement opportunities")
    pdf.bullet("Early Warning — At-risk trainee identification (86% accuracy)")
    pdf.bullet("Smart Recommendations — Actionable suggestions")
    pdf.bullet("Auto Reports — AI-generated comprehensive reports")

    pdf.section_heading("5. Technical Specifications")
    pdf.simple_table(
        ["Feature", "Details"],
        [
            ["Pages", "22 pages"],
            ["API Endpoints", "39+ routes"],
            ["Database", "44 tables (PostgreSQL)"],
            ["Authentication", "JWT with 59 permissions"],
            ["Hosting", "Cloud (Vercel + Neon)"],
            ["Languages", "Arabic (RTL) + English (LTR)"],
            ["Reports", "PDF + Excel export"],
        ]
    )

    pdf.section_heading("6. Current Status")
    pdf.simple_table(
        ["Item", "Status"],
        [
            ["Development", "Complete (v1.0)"],
            ["Hosting", "Active on Vercel"],
            ["Database", "Active on Neon PostgreSQL"],
            ["Demo Data", "Loaded (377 users)"],
            ["AI Services", "Active and working"],
            ["URL", "campus27.vercel.app"],
            ["Demo Login", "admin@campus27.sa / 123456"],
        ]
    )

    pdf.section_heading("7. Success Criteria")
    pdf.bullet("80% user adoption within 30 days of pilot launch")
    pdf.bullet("40% improvement in administrative efficiency")
    pdf.bullet("100% quality documentation ready for accreditation")
    pdf.bullet("User satisfaction >= 4/5")

    pdf.output(OUT / "Campus27-BRD.pdf")
    print("✓ Campus27-BRD.pdf")


def generate_pilot_pdf():
    pdf = ArabicPDF()
    pdf.add_page()

    pdf.header_block("Campus27", "Pilot Launch Plan")
    pdf.body_text("Duration: 8 Weeks | Participants: ~30 People")

    pdf.section_heading("1. Pilot Objectives")
    pdf.simple_table(
        ["#", "Objective", "KPI", "Target"],
        [
            ["1", "Platform fitness", "Task completion rate", ">= 80%"],
            ["2", "Usability", "User satisfaction", ">= 4/5"],
            ["3", "Quality system test", "KPIs recorded", ">= 10"],
            ["4", "AI usefulness", "Prediction accuracy", ">= 80%"],
            ["5", "Feedback collection", "Issues resolved", ">= 90%"],
        ]
    )

    pdf.section_heading("2. Participants")
    pdf.simple_table(
        ["Category", "Count", "Selection Criteria"],
        [
            ["College VPs", "3", "VP Trainers + VP Trainees + VP Quality"],
            ["Dept Heads", "3", "Diverse departments"],
            ["Trainers", "20", "6-7 per participating department"],
            ["Quality Officer", "1", "Quality module testing"],
            ["Accountant", "1", "Finance module testing"],
            ["IT + Coordinator", "2", "Support and coordination"],
        ]
    )

    pdf.section_heading("3. Timeline")

    pdf.body_text("Phase 1: Setup (Week 1)")
    pdf.bullet("Create real user accounts")
    pdf.bullet("Enter actual department and course data")
    pdf.bullet("Enter participating trainer data")
    pdf.bullet("Comprehensive technical testing")

    pdf.body_text("Phase 2: Training (Week 2)")
    pdf.simple_table(
        ["Session", "Audience", "Duration", "Content"],
        [
            ["1", "College VPs", "2 hours", "Dashboard, Reports, AI"],
            ["2", "Dept Heads", "3 hours", "Dept management, courses"],
            ["3-4", "Trainers", "3 hrs x 2", "Grades, attendance, tasks"],
            ["5", "Quality Officer", "2 hours", "Full quality system"],
        ]
    )

    pdf.body_text("Phase 3: Active Trial (Weeks 3-6)")
    pdf.simple_table(
        ["Week", "Focus", "Activities"],
        [
            ["3", "Soft start", "Daily attendance + one exam grades"],
            ["4", "Expansion", "Tasks + finance + quality KPIs"],
            ["5", "Stabilization", "All daily operations via platform"],
            ["6", "Initial review", "Review reports + evaluate AI"],
        ]
    )

    pdf.body_text("Phase 4: Evaluation & Decision (Weeks 7-8)")
    pdf.bullet("User satisfaction survey")
    pdf.bullet("Collect and analyze feedback")
    pdf.bullet("Pilot results report")
    pdf.bullet("Present results to Dean")
    pdf.bullet("Expansion or adjustment decision")

    pdf.section_heading("4. Success Criteria")
    pdf.simple_table(
        ["Criterion", "Minimum", "Ideal"],
        [
            ["Daily usage rate", "60%", "80%+"],
            ["User satisfaction", "3.5/5", "4+/5"],
            ["Critical issues", "<= 3", "0"],
            ["Task completion time", "Faster than current", "50%+ faster"],
        ]
    )

    pdf.section_heading("5. Next Steps")
    pdf.bullet("Dean approval of pilot plan")
    pdf.bullet("Select 3 participating departments")
    pdf.bullet("Assign project coordinator")
    pdf.bullet("Proposed start: 2 weeks after approval")

    pdf.output(OUT / "Campus27-Pilot-Plan.pdf")
    print("✓ Campus27-Pilot-Plan.pdf")


# ═══════════════════════════════════════════════════════════════
# MAIN — Generate all documents
# ═══════════════════════════════════════════════════════════════
if __name__ == "__main__":
    print("Generating Campus27 launch package...\n")

    print("── Word Documents (DOCX) ──")
    generate_brd()
    generate_pilot_plan()

    print("\n── PowerPoint (PPTX) ──")
    generate_pptx()

    print("\n── PDF Documents ──")
    generate_brd_pdf()
    generate_pilot_pdf()

    print(f"\n✅ All files generated in: {OUT}")
    print("\nFiles:")
    for f in sorted(OUT.glob("*.docx")) + sorted(OUT.glob("*.pptx")) + sorted(OUT.glob("*.pdf")):
        size_kb = f.stat().st_size / 1024
        print(f"  {f.name:40s} ({size_kb:.0f} KB)")
